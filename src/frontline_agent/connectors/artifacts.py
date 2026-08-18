"""Artifact generation and delivery.

R3.2 names "documents, prototypes, presentations" and is the requirement that
went longest unaddressed. The obvious answer was a Google Workspace or M365
connector; the better answer is that the agent already writes well, so the only
thing actually missing was a way to put what it wrote somewhere useful.

Two destinations, chosen by shape:

  Canvas   anything text-shaped. Lives in Slack, editable by him, shareable with
           a link, readable on a phone with no download. For "draft the spec"
           this beats a .docx by a wide margin — a downloaded file on a phone is
           a dead end.

  File     anything binary. Decks, exports, images. Uploaded into the thread
           where he asked for it.

Both are write_internal: they create something inside Frontline, and a human can
delete it. Neither leaves the building, which is why neither is write_external
even though a canvas can be shared onward afterward.
"""

from __future__ import annotations

import io
from typing import Any

from ..config import get_settings
from .base import Connector, DataClass, Risk

connector = Connector(
    name="artifacts",
    description="Create Slack canvases and upload generated files.",
    provider="slack-user",
)


def _client(token: str | None):
    from slack_sdk.web.async_client import AsyncWebClient

    settings = get_settings()
    return AsyncWebClient(token=token or settings.slack_bot_token)


@connector.tool(
    "create_canvas",
    "Create a Slack canvas from markdown. Use for specs, briefs, summaries, "
    "meeting notes, strategy drafts — anything the principal will read or edit "
    "rather than file away. Prefer this over uploading a document.",
    {
        "type": "object",
        "properties": {
            "title": {"type": "string"},
            "markdown": {
                "type": "string",
                "description": "Canvas body. Supports headings, lists, tables, and links.",
            },
            "channel": {
                "type": "string",
                "description": "Optional channel to share the canvas into after creating it.",
            },
        },
        "required": ["title", "markdown"],
    },
    risk=Risk.WRITE_INTERNAL,
    data_class=DataClass.INTERNAL,
)
async def create_canvas(args: dict[str, Any], *, token: str | None) -> Any:
    if get_settings().is_local:
        return {
            "ok": True,
            "simulated": True,
            "canvas_id": "F_LOCAL_CANVAS",
            "title": args["title"],
            "chars": len(args["markdown"]),
        }

    client = _client(token)
    response = await client.canvases_create(
        title=args["title"],
        document_content={"type": "markdown", "markdown": args["markdown"]},
    )
    canvas_id = response.get("canvas_id")

    if channel := args.get("channel"):
        await client.canvases_access_set(
            canvas_id=canvas_id, access_level="write", channel_ids=[channel]
        )
    return {"ok": True, "canvas_id": canvas_id, "title": args["title"]}


@connector.tool(
    "upload_file",
    "Upload a text file into a Slack conversation. Use for CSVs, exports, and "
    "anything the principal needs outside Slack. For prose, use create_canvas instead.",
    {
        "type": "object",
        "properties": {
            "filename": {"type": "string"},
            "content": {"type": "string"},
            "title": {"type": "string"},
            "channel": {"type": "string"},
            "thread_ts": {"type": "string"},
        },
        "required": ["filename", "content", "channel"],
    },
    risk=Risk.WRITE_INTERNAL,
    data_class=DataClass.INTERNAL,
)
async def upload_file(args: dict[str, Any], *, token: str | None) -> Any:
    if get_settings().is_local:
        return {
            "ok": True,
            "simulated": True,
            "filename": args["filename"],
            "bytes": len(args["content"].encode()),
        }

    response = await _client(token).files_upload_v2(
        filename=args["filename"],
        content=args["content"],
        title=args.get("title", args["filename"]),
        channel=args["channel"],
        thread_ts=args.get("thread_ts"),
    )
    return {"ok": True, "file_id": response.get("file", {}).get("id")}


@connector.tool(
    "render_deck",
    "Render a slide deck to PowerPoint and upload it. Give one slide object per "
    "slide with a title and bullet points; write the content yourself rather than "
    "asking the principal to outline it.",
    {
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "Deck title, used on the title slide."},
            "subtitle": {"type": "string"},
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {"type": "string"},
                        "bullets": {"type": "array", "items": {"type": "string"}},
                    },
                    "required": ["heading"],
                },
            },
            "channel": {"type": "string"},
            "thread_ts": {"type": "string"},
        },
        "required": ["title", "slides", "channel"],
    },
    risk=Risk.WRITE_INTERNAL,
    data_class=DataClass.INTERNAL,
)
async def render_deck(args: dict[str, Any], *, token: str | None) -> Any:
    try:
        payload = _build_pptx(args)
    except ImportError:
        return {
            "ok": False,
            "error": "python-pptx is not installed. Install the 'artifacts' extra: "
            "pip install -e '.[artifacts]'",
        }

    filename = f"{args['title'].replace(' ', '-').lower()}.pptx"

    if get_settings().is_local:
        return {
            "ok": True,
            "simulated": True,
            "filename": filename,
            "slides": len(args["slides"]) + 1,
            "bytes": len(payload),
        }

    response = await _client(token).files_upload_v2(
        filename=filename,
        file=payload,
        title=args["title"],
        channel=args["channel"],
        thread_ts=args.get("thread_ts"),
    )
    return {"ok": True, "file_id": response.get("file", {}).get("id"), "slides": len(args["slides"]) + 1}


def _build_pptx(args: dict[str, Any]) -> bytes:
    """Structured outline -> PPTX bytes.

    Deliberately plain: title slide, then title-and-content per slide, using the
    default template. Frontline has brand templates and this does not know about
    them; pointing `Presentation()` at a .potx is the whole fix once someone
    supplies one. Producing something unbranded and obviously draft is better
    than producing something that looks official and isn't.
    """
    from pptx import Presentation
    from pptx.util import Pt

    deck = Presentation()

    title_slide = deck.slides.add_slide(deck.slide_layouts[0])
    title_slide.shapes.title.text = args["title"]
    if subtitle := args.get("subtitle"):
        title_slide.placeholders[1].text = subtitle

    for slide_spec in args["slides"]:
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = slide_spec["heading"]
        bullets = slide_spec.get("bullets") or []
        if not bullets:
            continue
        frame = slide.placeholders[1].text_frame
        frame.text = bullets[0]
        for bullet in bullets[1:]:
            paragraph = frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
        for paragraph in frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)

    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()
